#!/usr/bin/env python3
"""Build Allada deck with python-pptx (PowerPoint-native OOXML)."""
import json
import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).parent
DATA = json.loads((ROOT / "data.json").read_text(encoding="utf-8"))
OUT = ROOT.parent / "Allada_Stocker_Tables_Overview.pptx"
TOTAL = 16

W = Inches(13.333)
H = Inches(7.5)

INK = RGBColor(0x15, 0x23, 0x3B)
TEAL = RGBColor(0x0E, 0x7C, 0x7B)
TEALL = RGBColor(0x8F, 0xCF, 0xCD)
AMBER = RGBColor(0xE0, 0xA1, 0x3A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG = RGBColor(0xF4, 0xF6, 0xF8)
TX = RGBColor(0x1F, 0x2A, 0x37)
MUT = RGBColor(0x66, 0x73, 0x85)
LINE = RGBColor(0xE2, 0xE8, 0xF0)
PALE = RGBColor(0xEE, 0xF2, 0xF5)
LIGHT = RGBColor(0xC7, 0xD2, 0xDE)

BC = [
    RGBColor(0x0E, 0x7C, 0x7B),
    RGBColor(0x1C, 0x6F, 0xB0),
    RGBColor(0x6B, 0x4F, 0xA0),
    RGBColor(0xC2, 0x79, 0x2B),
    RGBColor(0xB2, 0x41, 0x5A),
    RGBColor(0x4F, 0x7A, 0x2E),
]

BUCKETS = [
    ("Bloomington", "UAS / RNAi", "Stock-center, overexpression / knockdown"),
    ("Vienna (VDRC)", "UAS / RNAi", "VDRC knockdown collection"),
    ("Bloomington", "Allele / Insertion", "Stock-center classical alleles"),
    ("Stock Center", "Allele / Insertion", "Other orderable centers (Kyoto, etc.)"),
    ("Non-Stock-Center", "UAS / RNAi", "Custom phenotype reagents (no FBst)"),
    ("Non-Stock-Center", "Allele / Insertion", "Custom phenotype reagents (no FBst)"),
]

OVERVIEW_ROWS = [
    ("vGAT screen", "GABA gene list", "140", "Bulk-seq vGAT hits (Aldeb)", "Curated list (5 col)"),
    ("Priority v2", "BPD-GWAS orthologs", "168", "Bipolar GWAS, Nature 2025 S31", "DIOPT ortholog (19 col)"),
    ("", "DIOPT Bellenguez/Lambert", "227", "Alzheimer's GWAS orthologs", "DIOPT + stocks (36 col)"),
    ("", "Jordan VCM", "35", "Curated candidate genes", "Simple list (4 col)"),
    ("", "Slumber neuropeptides", "33", "Curated neuropeptides/transmitters", "Simple list (4 col)"),
    ("", "Slumber RNAi candidates", "16", "Curated RNAi candidates", "Simple list (4 col)"),
    ("Transcriptomics", "CSW FC0.5 (7 tiers)", "788", "Sleep/wake cycling + correlation", "Cycling + AI refs (31 col)"),
    ("", "CSW FC1 (4 tiers)", "114", "Sleep/wake cycling (stricter FC)", "Cycling + AI refs (31 col)"),
    ("", "SleepHistory correlates (2)", "215", "History +/- literature correlates", "Lit. category (18 col)"),
    ("", "History x Rebound overlaps (2)", "20", "Overlapping history/rebound sets", "Lit. category (14 col)"),
    ("", "Homeostatic genes (1)", "6", "Sleep homeostasis set", "Simple list (4 col)"),
]

TX_NM = {
    "CSW FC0.5": "CSW FC0.5",
    "CSW FC1": "CSW FC1",
    "Homeostatic Genes": "Homeostatic genes",
    "Overlapping SleepHistory<>SleepRebound Correlates": "History x Rebound overlaps",
    "SleepHistory Correlates": "SleepHistory correlates",
}

TARGETED_ORDER = [
    "vGAT — GABA gene list",
    "BPD-GWAS orthologs",
    "DIOPT Bellenguez/Lambert (AD)",
    "Jordan VCM",
    "Slumber Neuropeptides+transmitters",
    "Slumber RNAi candidates",
]

TARGETED_LABELS = {
    "vGAT — GABA gene list": "vGAT — GABA gene list",
    "BPD-GWAS orthologs": "BPD-GWAS orthologs",
    "DIOPT Bellenguez/Lambert (AD)": "DIOPT Bellenguez/Lambert (AD)",
    "Jordan VCM": "Jordan VCM",
    "Slumber Neuropeptides+transmitters": "Slumber neuropeptides",
    "Slumber RNAi candidates": "Slumber RNAi candidates",
}

BAR_LABELS = ["GABA", "BPD-GWAS", "DIOPT AD", "Jordan", "Slumber NP", "Slumber RNAi"]
BUCKET_SHORT = ["Bloom · UAS", "Vienna · UAS", "Bloom · Allele", "Stock Ctr · Allele", "Non-SC · UAS", "Non-SC · Allele"]

SOURCE_GROUPS = [
    ("1", "Input gene sets", "22 tables: screen hits, human GWAS orthologs, curated candidates, transcriptomics tiers"),
    ("2", "FlyBase stock expansion", "For each FBgn, collect stocks, alleles, constructs, genotypes, phenotypes, and references"),
    ("3", "Phenotype + literature gates", "Keep phenotype-backed reagents; flag sleep/circadian/rhythm/locomotor support"),
    ("4", "Priority sorting", "Assign each stock once into six orderability/reagent buckets"),
    ("5", "Review-ready workbook", "Six bucket sheets + references + all phenotypic stocks + contents/counts"),
]

TARGETED_PROVENANCE = [
    (
        "vGAT screen",
        "Aldeb bulk-seq vGAT hits",
        "Internal screen-derived gene list; bottom-10% vGAT hits, supplied as symbols + FBgn IDs.",
        "140 fly genes directly stocked.",
    ),
    (
        "Bipolar GWAS",
        "O'Connell, Koromina, van der Veen et al.; published 22 Jan 2025; Nature 639:968-975.",
        "Human bipolar-disorder GWAS: 158,036 BD cases + 2.8M controls; 298 loci. Supplementary Table S31 was mapped through DIOPT and Filter2.",
        "168 fly ortholog rows; DIOPT score/rank retained.",
    ),
    (
        "Alzheimer's GWAS",
        "Bellenguez/Lambert AD GWAS-derived ortholog list.",
        "Human Alzheimer's disease candidate genes mapped to Drosophila orthologs and paired with pre-pulled BDSC RNAi hints.",
        "227 fly genes with DIOPT confidence + stock hints.",
    ),
    (
        "Jordan VCM",
        "Internal curated candidate list.",
        "Focused VCM candidate genes supplied as current symbols + FBgn IDs.",
        "35 fly genes directly stocked.",
    ),
    (
        "Slumber NP/transmitters",
        "Internal curated Slumber list.",
        "Neuropeptide and transmitter genes likely relevant to sleep/wake phenotypes.",
        "33 fly genes directly stocked.",
    ),
    (
        "Slumber RNAi",
        "Internal curated Slumber list.",
        "Short RNAi candidate list for targeted follow-up.",
        "16 fly genes directly stocked.",
    ),
]

TX_REFERENCE = (
    "Rosensweig, Shah, Sisobhan, Andreani & Allada",
    "Shared Transcriptomic Responses to Distinct Sleep-Wake Manipulations Reveals Multiple Homeostatic Pathways in Drosophila",
    "posted 6 Jun 2026; bioRxiv preprint / lab manuscript submitted for review; doi: 10.64898/2026.02.28.708752",
)

TX_DATASETS = [
    ("Baseline time course", "Whole-brain RNA-seq at ZT0/4/8/12/16/20; validates clock transcript cycling."),
    ("Mechanical SD", "3h, 6h, and 12h sleep deprivation ending at ZT0; each followed by rebound assessment."),
    ("Thermogenetic sleep", "R85C10-GAL4 > TrpA1 activation promoted sleep and subsequent anti-rebound."),
    ("Published induction sets", "Optogenetic R23E1040 > ChRimson and pharmacologic THIP whole-brain RNA-seq."),
    ("Correlation extension", "Sleep-history/rebound correlations across 9 datasets; added WakeGAL4 TrpA1 lines and excluded opto/THIP where direct sleep/rebound metadata was unavailable."),
]

PIPELINE_EXAMPLE = [
    ("Human gene", "MED24 from Nature 2025 bipolar GWAS"),
    ("Orthology", "DIOPT maps MED24 to fly MED24 (FBgn0035851), high confidence"),
    ("Stock expansion", "Find FlyBase stocks/alleles/constructs tied to MED24"),
    ("Priority placement", "Place orderable phenotype-backed reagents before custom/non-stock-center reagents"),
]


def short_bucket(combo: str) -> str:
    return (
        combo.replace(" >> ", " · ")
        .replace("AlleleOrInsertion", "Allele/Insertion")
        .replace("Non Stock Center", "Non-Stock-Center")
        .replace("Vienna", "Vienna (VDRC)")
        .replace(" · Phenotype", "")
    )


def short_tier(name: str) -> str:
    base = name.split("_n=")[0]
    base = base.split(" Table")[0]
    return base.replace("_", " ").replace("frequency", "freq").replace("Tx-Omics ", "")


def tier_rows(keys):
    out = []
    for k in keys:
        for it in DATA["tx_tiers"].get(k, []):
            out.append([short_tier(it[0]), str(it[1])])
    return out


def new_prs():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def fill_bg(slide, color=BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def set_run(p, text, *, size=11, bold=False, color=TX, font="Calibri", align=None):
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.name = font
    p.font.color.rgb = color
    if align:
        p.alignment = align


def textbox(slide, left, top, width, height, text, *, size=11, bold=False, color=TX, font="Calibri", align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    set_run(tf.paragraphs[0], text, size=size, bold=bold, color=color, font=font, align=align)
    return tb


def rect(slide, left, top, width, height, fill, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line:
        sh.line.color.rgb = line
        sh.line.width = Pt(0.75)
    else:
        sh.line.fill.background()
    return sh


def header(slide, kicker, title):
    rect(slide, Inches(0), Inches(0), W, Inches(1.18), INK)
    rect(slide, Inches(0), Inches(1.18), W, Inches(0.06), AMBER)
    if kicker:
        textbox(slide, Inches(0.55), Inches(0.18), Inches(12), Inches(0.3), kicker.upper(), size=11, bold=True, color=TEALL)
    textbox(slide, Inches(0.55), Inches(0.44), Inches(12.2), Inches(0.62), title, size=25, bold=True, color=WHITE, font="Georgia")


def footer(slide, n):
    rect(slide, Inches(0), Inches(7.18), W, Inches(0.32), INK)
    textbox(slide, Inches(0.5), Inches(7.16), Inches(8), Inches(0.32), "Allada Lab  ·  FlyBase Reagent Stocker", size=9, color=TEALL)
    textbox(slide, Inches(11.8), Inches(7.16), Inches(1.2), Inches(0.32), f"{n} / {TOTAL}", size=9, color=TEALL, align=PP_ALIGN.RIGHT)


def style_cell(cell, text, *, bold=False, size=10, color=TX, fill=None, align=PP_ALIGN.LEFT):
    cell.text = text
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if fill:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.name = "Calibri"
        p.font.color.rgb = color
        p.alignment = align


def add_table(slide, left, top, width, headers, rows, col_widths, row_size=10, header_size=11):
    nrows, ncols = len(rows) + 1, len(headers)
    th = slide.shapes.add_table(nrows, ncols, left, top, width, Inches(0.32 * nrows)).table
    for i, w in enumerate(col_widths):
        th.columns[i].width = w
    for j, h in enumerate(headers):
        style_cell(th.cell(0, j), h, bold=True, size=header_size, color=WHITE, fill=INK)
    for i, row in enumerate(rows):
        bg = PALE if i % 2 else WHITE
        for j, val in enumerate(row):
            bold = isinstance(val, dict)
            txt = val.get("t", "") if bold else str(val)
            col = val.get("c", TX) if bold else TX
            style_cell(th.cell(i + 1, j), txt, bold=bold, size=row_size, color=col, fill=bg)
    return th


def stat_card(slide, left, top, value, label, color):
    rect(slide, left, top, Inches(2.92), Inches(1.05), WHITE, LINE)
    textbox(slide, left, top + Inches(0.12), Inches(2.92), Inches(0.62), str(value), size=30, bold=True, color=color, font="Georgia", align=PP_ALIGN.CENTER)
    textbox(slide, left, top + Inches(0.74), Inches(2.92), Inches(0.34), label, size=10.5, color=MUT, align=PP_ALIGN.CENTER)


def callout(slide, left, top, width, height, title, body, color=TEAL, body_size=10):
    rect(slide, left, top, Inches(0.12), height, color)
    rect(slide, left + Inches(0.12), top, width - Inches(0.12), height, WHITE, LINE)
    textbox(slide, left + Inches(0.28), top + Inches(0.12), width - Inches(0.42), Inches(0.28), title, size=12, bold=True, color=INK)
    textbox(slide, left + Inches(0.28), top + Inches(0.45), width - Inches(0.42), height - Inches(0.54), body, size=body_size, color=MUT)


def bullet_rows(slide, left, top, width, items, *, size=10.5, gap=0.42, color=TEAL):
    y = top
    for title, body in items:
        oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, y + Inches(0.06), Inches(0.13), Inches(0.13))
        oval.fill.solid()
        oval.fill.fore_color.rgb = color
        oval.line.fill.background()
        textbox(slide, left + Inches(0.25), y, width - Inches(0.25), Inches(0.22), title, size=size, bold=True, color=TX)
        textbox(slide, left + Inches(0.25), y + Inches(0.22), width - Inches(0.25), Inches(0.28), body, size=size - 1.2, color=MUT)
        y += Inches(gap)


def numbered_flow(slide, left, top, width, rows, *, color=TEAL):
    box_w = int(width / len(rows))
    for i, (num, title, body) in enumerate(rows):
        x = left + box_w * i
        rect(slide, x, top, box_w - Inches(0.12), Inches(1.45), WHITE, LINE)
        oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.16), top + Inches(0.2), Inches(0.42), Inches(0.42))
        oval.fill.solid()
        oval.fill.fore_color.rgb = color
        oval.line.fill.background()
        textbox(slide, x + Inches(0.16), top + Inches(0.2), Inches(0.42), Inches(0.42), num, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        textbox(slide, x + Inches(0.68), top + Inches(0.18), box_w - Inches(0.92), Inches(0.3), title, size=11.5, bold=True, color=INK)
        textbox(slide, x + Inches(0.18), top + Inches(0.72), box_w - Inches(0.42), Inches(0.56), body, size=8.5, color=MUT)


def hbars(slide, left, top, width, rows, color, rh=Inches(0.5), label_w=Inches(1.5), val_w=Inches(0.85)):
    max_v = max(r[1] for r in rows) or 1
    bar_max = width - label_w - val_w - Inches(0.12)
    for i, (lab, val) in enumerate(rows):
        y = top + rh * i
        textbox(slide, left, y, label_w - Inches(0.1), rh, lab, size=9.5, color=MUT, align=PP_ALIGN.RIGHT)
        bw = max(Inches(0.05), int(bar_max * val / max_v))
        rect(slide, left + label_w, y + rh * 0.2, bw, rh * 0.6, color)
        textbox(slide, left + label_w + bw + Inches(0.06), y, val_w, rh, f"{val:,}", size=9.5, bold=True, color=TX)


def bucket_table(slide, left, top, width, buckets):
    headers = ["Priority bucket", "Stocks", "Alleles", "Genes"]
    rows = [[short_bucket(b[0]), b[1], b[2], b[3]] for b in buckets]
    add_table(slide, left, top, width, headers, rows, [Inches(3.2), Inches(0.9), Inches(0.9), Inches(0.85)])


def dataset_slide(prs, n, kicker, title, n_genes, total, buckets, schema_note, cols, foot=None):
    slide = blank_slide(prs)
    fill_bg(slide)
    header(slide, kicker, title)
    stats = [
        ("Input genes", n_genes, AMBER),
        ("Unique stocks", total[1], TEAL),
        ("Unique alleles", total[2], TEAL),
        ("Genes w/ stocks", total[3], TEAL),
    ]
    x = Inches(0.55)
    for lab, val, col in stats:
        stat_card(slide, x, Inches(1.5), val, lab, col)
        x += Inches(3.05)
    rect(slide, Inches(0.55), Inches(2.78), Inches(5.55), Inches(3.7), WHITE, LINE)
    textbox(slide, Inches(0.78), Inches(2.92), Inches(5.1), Inches(0.34), "Input schema", size=13.5, bold=True, color=INK)
    textbox(slide, Inches(0.78), Inches(3.24), Inches(5.1), Inches(0.5), schema_note, size=9.5, color=MUT)
    y = Inches(3.78)
    for c0, c1 in cols:
        textbox(slide, Inches(1.02), y, Inches(5), Inches(0.32), f"{c0} — {c1}", size=9.8, color=TX)
        y += Inches(0.42)
    rect(slide, Inches(6.3), Inches(2.78), Inches(6.5), Inches(3.7), WHITE, LINE)
    textbox(slide, Inches(6.53), Inches(2.92), Inches(6), Inches(0.34), "Stocks per priority bucket", size=13.5, bold=True, color=INK)
    bucket_table(slide, Inches(6.53), Inches(3.34), Inches(5.85), buckets)
    if foot:
        textbox(slide, Inches(6.53), Inches(6.12), Inches(6), Inches(0.32), foot, size=9, color=MUT)
    footer(slide, n)


def tx_slide(prs, n, kicker, title, schema_note, cols, right_title, tiers, foot=None):
    slide = blank_slide(prs)
    fill_bg(slide)
    header(slide, kicker, title)
    rect(slide, Inches(0.55), Inches(1.5), Inches(5.55), Inches(4.98), WHITE, LINE)
    textbox(slide, Inches(0.78), Inches(1.64), Inches(5.1), Inches(0.34), "Input schema", size=13.5, bold=True, color=INK)
    textbox(slide, Inches(0.78), Inches(1.97), Inches(5.15), Inches(0.62), schema_note, size=9.5, color=MUT)
    y = Inches(2.66)
    for c0, c1 in cols:
        textbox(slide, Inches(1.02), y, Inches(5), Inches(0.34), f"{c0} — {c1}", size=9.8, color=TX)
        y += Inches(0.46)
    rect(slide, Inches(6.3), Inches(1.5), Inches(6.5), Inches(4.98), WHITE, LINE)
    textbox(slide, Inches(6.53), Inches(1.64), Inches(6), Inches(0.34), right_title, size=13.5, bold=True, color=INK)
    add_table(slide, Inches(6.53), Inches(2.04), Inches(5.85), ["Gene-set table (tier)", "Genes"], tiers, [Inches(4.85), Inches(1.0)])
    if foot:
        textbox(slide, Inches(6.53), Inches(6.05), Inches(6), Inches(0.36), foot, size=9, color=MUT)
    footer(slide, n)


def slide_title(prs):
    slide = blank_slide(prs)
    fill_bg(slide, INK)
    rect(slide, Inches(0), Inches(0), Inches(0.28), H, TEAL)
    rect(slide, Inches(0.28), Inches(0), Inches(0.08), H, AMBER)
    textbox(slide, Inches(1), Inches(1.6), Inches(11), Inches(0.4), "FLYBASE REAGENT STOCKER", size=14, bold=True, color=TEALL)
    textbox(slide, Inches(0.95), Inches(2.05), Inches(11.6), Inches(1), "From Gene Sets to Fly Stocks", size=43, bold=True, color=WHITE, font="Georgia")
    textbox(slide, Inches(1), Inches(3.15), Inches(11.4), Inches(0.5), "End-to-end prioritization of Drosophila reagents for sleep/circadian follow-up", size=18, color=LIGHT)
    items = [("22", "input gene-set tables"), ("3", "source families"), ("6", "priority buckets"), ("8,300+", "stocks organized")]
    x = Inches(1)
    for val, lab in items:
        textbox(slide, x, Inches(4.5), Inches(2.7), Inches(0.6), val, size=30, bold=True, color=AMBER, font="Georgia")
        textbox(slide, x, Inches(5.1), Inches(2.7), Inches(0.5), lab, size=12, color=LIGHT)
        x += Inches(2.85)
    textbox(slide, Inches(1), Inches(6.5), Inches(8), Inches(0.4), "Allada Lab · FlyBase Reagent Stocker", size=13, color=TEALL)


def slide_pipeline(prs):
    slide = blank_slide(prs)
    fill_bg(slide)
    header(slide, "End-to-end pipeline", "What happens to each input gene set")
    numbered_flow(slide, Inches(0.55), Inches(1.48), Inches(12.35), SOURCE_GROUPS)
    textbox(slide, Inches(0.55), Inches(3.28), Inches(12), Inches(0.34), "Priority buckets: orderability first, then reagent class; each stock appears once", size=13, bold=True, color=INK)
    bw = Inches(4.0)
    for i, (b0, b1, b2) in enumerate(BUCKETS):
        col, row = i % 3, i // 3
        X = Inches(0.55) + col * (bw + Inches(0.18))
        Y = Inches(3.72) + row * Inches(0.94)
        rect(slide, X, Y, Inches(0.13), Inches(0.78), BC[i])
        rect(slide, X + Inches(0.13), Y, bw - Inches(0.13), Inches(0.78), WHITE, LINE)
        textbox(slide, X + Inches(0.28), Y + Inches(0.08), bw - Inches(0.45), Inches(0.32), f"{b0}  ·  {b1}", size=11.5, bold=True, color=TX)
        textbox(slide, X + Inches(0.28), Y + Inches(0.4), bw - Inches(0.45), Inches(0.3), b2, size=9, color=MUT)
    textbox(slide, Inches(0.55), Inches(6.45), Inches(12.2), Inches(0.5), "The practical question this answers: for any candidate gene, which phenotype-backed fly reagents are easiest to order and most relevant to sleep/circadian biology?", size=10.5, color=MUT)
    footer(slide, 3)


def slide_executive_summary(prs):
    slide = blank_slide(prs)
    fill_bg(slide)
    header(slide, "Overview", "What this stocker run produces")
    stat_card(slide, Inches(0.55), Inches(1.5), "22", "input gene-set tables", AMBER)
    stat_card(slide, Inches(3.6), Inches(1.5), "619", "targeted input genes", TEAL)
    stat_card(slide, Inches(6.65), Inches(1.5), str(DATA["tx_total_stocks"]), "transcriptomics stocks", TEAL)
    stat_card(slide, Inches(9.7), Inches(1.5), "6", "priority buckets per set", AMBER)

    callout(slide, Inches(0.55), Inches(3.05), Inches(5.95), Inches(1.35), "Core message", "The pipeline turns heterogeneous biological evidence into orderable, phenotype-backed fly reagents, preserving source provenance and literature support so follow-up choices are auditable.", TEAL, body_size=10.5)
    callout(slide, Inches(6.85), Inches(3.05), Inches(5.95), Inches(1.35), "Narrative flow", "The deck starts with source provenance and prioritization logic, then moves into examples and deeper dataset-specific evidence.", AMBER, body_size=10.5)

    rows = [
        ("Targeted runs", "GABA/vGAT screen, bipolar GWAS orthologs, AD orthologs, and curated Slumber/Jordan lists."),
        ("Transcriptomics run", "Rosensweig/Shah sleep-wake manuscript-derived gene sets: CSW tiers, sleep-history/rebound correlations, and homeostatic candidates."),
        ("Output", "Every input table becomes a workbook with six priority sheets, references, all phenotypic stocks, and count summaries."),
    ]
    bullet_rows(slide, Inches(0.78), Inches(4.78), Inches(11.6), rows, gap=0.5)
    footer(slide, 2)


def slide_source_map(prs):
    slide = blank_slide(prs)
    fill_bg(slide)
    header(slide, "Source map", "Three source families feed the same prioritization workflow")
    cards = [
        ("Screen-derived", "GABA gene list from Aldeb bulk-seq vGAT hits (bottom 10%). Direct fly genes with FBgn IDs.", TEAL),
        ("Human genetics", "Bipolar and Alzheimer's disease human gene sets mapped to Drosophila through DIOPT, retaining orthology confidence.", AMBER),
        ("Lab/curated sleep biology", "Jordan VCM, Slumber neuropeptide/transmitter, Slumber RNAi, and Rosensweig/Shah transcriptomics-derived tiers.", BC[2]),
    ]
    x = Inches(0.55)
    for title, body, color in cards:
        callout(slide, x, Inches(1.55), Inches(4.0), Inches(1.65), title, body, color, body_size=10)
        x += Inches(4.18)
    textbox(slide, Inches(0.55), Inches(3.65), Inches(12), Inches(0.34), "Why the source split matters", size=14, bold=True, color=INK)
    rows = [
        ("Direct fly gene sets", "No orthology step; stocker starts from FBgn IDs and expands to available reagents."),
        ("Human disease gene sets", "Interpretation depends on orthology confidence and source study definition; DIOPT score/rank are therefore shown."),
        ("Transcriptomics tiers", "Genes are evidence-ranked by recurrence across sleep-wake manipulations or by history/rebound correlations."),
    ]
    bullet_rows(slide, Inches(0.8), Inches(4.12), Inches(11.5), rows, gap=0.58)
    textbox(slide, Inches(0.55), Inches(6.55), Inches(12.2), Inches(0.34), "All three source families converge to the same stock workbook format, so downstream comparison is apples-to-apples.", size=9.5, color=MUT)
    footer(slide, 4)


def slide_targeted_provenance(prs):
    slide = blank_slide(prs)
    fill_bg(slide)
    header(slide, "Dataset provenance", "Targeted gene-list inputs")
    rows = []
    for source, ref, desc, inp in TARGETED_PROVENANCE:
        rows.append([source, ref, desc, inp])
    add_table(
        slide,
        Inches(0.45),
        Inches(1.45),
        Inches(12.45),
        ["Dataset", "Reference", "Dataset description", "Stocker input"],
        rows,
        [Inches(1.55), Inches(3.25), Inches(4.9), Inches(2.75)],
        row_size=7.5,
        header_size=9,
    )
    textbox(slide, Inches(0.55), Inches(6.52), Inches(12), Inches(0.36), "Bipolar GWAS is now separated into reference metadata plus the downstream DIOPT-derived fly-gene table used by this run.", size=9.5, color=MUT)
    footer(slide, 5)


def slide_transcriptomics_provenance(prs):
    slide = blank_slide(prs)
    fill_bg(slide)
    header(slide, "Dataset provenance", "Rosensweig/Shah transcriptomics inputs")
    callout(slide, Inches(0.55), Inches(1.45), Inches(12.25), Inches(1.05), TX_REFERENCE[1], f"{TX_REFERENCE[0]}; {TX_REFERENCE[2]}", TEAL, body_size=10.5)
    textbox(slide, Inches(0.55), Inches(2.85), Inches(12), Inches(0.34), "How the manuscript datasets map into stocker inputs", size=13.5, bold=True, color=INK)
    y = Inches(3.25)
    for i, (title, body) in enumerate(TX_DATASETS):
        callout(slide, Inches(0.55), y, Inches(5.95), Inches(0.72), title, body, BC[i % len(BC)], body_size=8.2)
        y += Inches(0.82)
    callout(slide, Inches(6.85), Inches(3.25), Inches(5.95), Inches(2.52), "Staged evidence", "CSW FC1 and FC0.5 tiers ask: which genes respond consistently across independent sleep/wake manipulations? Correlation tiers ask: which genes track prior sleep history or predict rebound? fl.ai columns retain literature support.", AMBER, body_size=10.2)
    callout(slide, Inches(6.85), Inches(5.95), Inches(5.95), Inches(0.92), "Interpretation note", "Transcriptomics gene counts can recur across tiers; stock counts are the cleaner headline metric.", BC[2], body_size=9.2)
    footer(slide, 6)


def slide_counts_overview(prs):
    slide = blank_slide(prs)
    fill_bg(slide)
    header(slide, "What came out", "How many stocks were organized")
    targeted_rows = []
    st, sa = 0, 0
    for name in TARGETED_ORDER:
        d = DATA["targeted"][name]
        t = d["total"]
        st += int(t[1])
        sa += int(t[2])
        targeted_rows.append([TARGETED_LABELS[name], str(d["n_genes"]), t[1], t[2], t[3]])
    targeted_rows.append([{"t": "All targeted", "c": TX}, "619", {"t": str(st), "c": TX}, {"t": str(sa), "c": TX}, "-"])
    rect(slide, Inches(0.55), Inches(1.45), Inches(6.25), Inches(4.9), WHITE, LINE)
    textbox(slide, Inches(0.78), Inches(1.62), Inches(5.8), Inches(0.3), "Targeted gene-list runs", size=13, bold=True, color=INK)
    add_table(slide, Inches(0.78), Inches(2.0), Inches(5.8), ["Dataset", "Genes", "Stocks", "Alleles", "w/ stocks"], targeted_rows, [Inches(2.35), Inches(0.7), Inches(0.9), Inches(0.9), Inches(0.95)], row_size=8.7, header_size=9.5)
    rect(slide, Inches(7.05), Inches(1.45), Inches(5.75), Inches(4.9), WHITE, LINE)
    textbox(slide, Inches(7.28), Inches(1.62), Inches(5.3), Inches(0.3), "Transcriptomics run", size=13, bold=True, color=INK)
    tx_rows = [[TX_NM.get(r[0], r[0]), str(r[1]), str(r[2])] for r in DATA["tx_per_dataset"]]
    add_table(slide, Inches(7.28), Inches(2.0), Inches(5.28), ["Dataset category", "Tiers", "Stocks"], tx_rows, [Inches(3.35), Inches(0.85), Inches(1.08)], row_size=8.9, header_size=9.5)
    textbox(slide, Inches(7.28), Inches(4.45), Inches(5.3), Inches(0.32), "Priority bucket load (all transcriptomics tiers)", size=11.5, bold=True, color=INK)
    b_rows = [(BUCKET_SHORT[i], int(b[1])) for i, b in enumerate(DATA["tx_by_bucket"])]
    hbars(slide, Inches(7.25), Inches(4.85), Inches(5.15), b_rows, AMBER, rh=Inches(0.32), label_w=Inches(1.65))
    textbox(slide, Inches(0.55), Inches(6.55), Inches(12.2), Inches(0.34), "Counts are de-duplicated within each workbook; transcriptomics genes may repeat across evidence tiers.", size=9.5, color=MUT)
    footer(slide, 7)


def slide_med24_example(prs):
    slide = blank_slide(prs)
    fill_bg(slide)
    header(slide, "Worked example", "How a human GWAS gene becomes a fly stock shortlist")
    numbered_flow(slide, Inches(0.55), Inches(1.45), Inches(12.35), [(str(i + 1), t, b) for i, (t, b) in enumerate(PIPELINE_EXAMPLE)], color=AMBER)
    callout(slide, Inches(0.55), Inches(3.35), Inches(5.95), Inches(1.25), "Input row", "MED24 appears in the bipolar GWAS-derived DIOPT table. The source row retains human symbol, fly ortholog symbol, FBgn, DIOPT score, weighted score, rank, and support methods.", TEAL, body_size=9.5)
    callout(slide, Inches(6.85), Inches(3.35), Inches(5.95), Inches(1.25), "Interpretation", "The detailed evidence and provenance remain traceable while the practical output is a short, ordered reagent list.", BC[2], body_size=9.5)
    textbox(slide, Inches(0.55), Inches(5.05), Inches(12.1), Inches(0.3), "Representative GABA bucket examples still appear next, showing what the final stock-level rows look like.", size=12, bold=True, color=INK)
    footer(slide, 8)


def slide_overview(prs):
    slide = blank_slide(prs)
    fill_bg(slide)
    header(slide, "Inputs", "The input datasets at a glance")
    rows = []
    for run, ds, genes, src, schema in OVERVIEW_ROWS:
        r0 = {"t": run, "c": TEAL} if run else ""
        rows.append([r0, ds, genes, src, schema])
    add_table(slide, Inches(0.55), Inches(1.5), Inches(12.25), ["Run", "Dataset", "Genes", "Source / meaning", "Input schema"], rows,
              [Inches(1.6), Inches(3.3), Inches(0.9), Inches(4.0), Inches(2.45)])
    textbox(slide, Inches(0.55), Inches(6.55), Inches(12.2), Inches(0.4), "22 input tables in total. Gene counts shown as submitted; the transcriptomics run groups 16 frequency tiers under five dataset categories.", size=10, color=MUT)
    footer(slide, 3)


def slide_anatomy(prs, n=10):
    slide = blank_slide(prs)
    fill_bg(slide)
    header(slide, "Output", "Anatomy of every generated workbook")
    rect(slide, Inches(0.55), Inches(1.5), Inches(5.7), Inches(5.0), WHITE, LINE)
    textbox(slide, Inches(0.78), Inches(1.66), Inches(5.2), Inches(0.4), "What each workbook contains", size=14, bold=True, color=INK)
    sheets = [
        ("Sheet 1 – 6", "One per priority bucket (color-coded above)"),
        ("References", "Supporting publications per stock / allele"),
        ("All Phenotypic Stocks", "Full union of phenotype-backed stocks"),
        ("Contents", "Prioritization rules + per-sheet count breakdown"),
    ]
    y = Inches(2.2)
    for t, d in sheets:
        rect(slide, Inches(0.78), y + Inches(0.05), Inches(0.1), Inches(0.62), TEAL)
        textbox(slide, Inches(0.98), y, Inches(5.1), Inches(0.34), t, size=12.5, bold=True, color=TX)
        textbox(slide, Inches(0.98), y + Inches(0.32), Inches(5.1), Inches(0.34), d, size=10, color=MUT)
        y += Inches(0.82)
    textbox(slide, Inches(0.78), Inches(5.7), Inches(5.25), Inches(0.7), "Priority rule: a stock is placed in the first bucket it qualifies for, so the six tables never double-count a stock.", size=10.5, color=TX)
    rect(slide, Inches(6.45), Inches(1.5), Inches(6.35), Inches(5.0), WHITE, LINE)
    textbox(slide, Inches(6.68), Inches(1.66), Inches(5.9), Inches(0.4), "Key columns (approx. 59 total), grouped", size=14, bold=True, color=INK)
    groups = [
        ("Identity", "FBst, stock_number, collection, genotype"),
        ("Reagent type", "RNAi, UAS, GAL4, mutant, sgRNA, custom_stock"),
        ("Gene / allele", "gene_symbols, flybase_gene_ids, AlleleSymbol, class"),
        ("Construct parts", "FBti / FBtp / FBal symbols + IDs, Balancers"),
        ("Literature", "PMID, total_refs, sleep/circadian ref counts"),
        ("Phenotype + similarity", "PHENOTYPE_RELEVANCE_SCORE, cosine similarities"),
    ]
    y = Inches(2.2)
    for i, (t, d) in enumerate(groups):
        oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.68), y + Inches(0.02), Inches(0.16), Inches(0.16))
        oval.fill.solid()
        oval.fill.fore_color.rgb = BC[i]
        oval.line.fill.background()
        textbox(slide, Inches(6.95), y, Inches(5.6), Inches(0.3), t, size=12, bold=True, color=TX)
        textbox(slide, Inches(6.95), y + Inches(0.26), Inches(5.7), Inches(0.3), d, size=9.5, color=MUT)
        y += Inches(0.72)
    footer(slide, n)


def slide_summary_targeted(prs):
    slide = blank_slide(prs)
    fill_bg(slide)
    header(slide, "Summary counts", "Targeted gene-list runs — totals")
    rows = []
    st, sa = 0, 0
    for name in TARGETED_ORDER:
        d = DATA["targeted"][name]
        t = d["total"]
        st += int(t[1])
        sa += int(t[2])
        rows.append([TARGETED_LABELS[name], str(d["n_genes"]), t[1], t[2], t[3]])
    rows.append([{"t": "All targeted runs", "c": TX}, "619", {"t": str(st), "c": TX}, {"t": str(sa), "c": TX}, "—"])
    rect(slide, Inches(0.55), Inches(1.5), Inches(6.55), Inches(4.55), WHITE, LINE)
    textbox(slide, Inches(0.78), Inches(1.64), Inches(6.1), Inches(0.34), "Unique stocks / alleles per dataset", size=13, bold=True, color=INK)
    add_table(slide, Inches(0.78), Inches(2.04), Inches(6.1), ["Dataset", "Genes", "Stocks", "Alleles", "w/ stocks"], rows,
              [Inches(2.5), Inches(0.78), Inches(0.95), Inches(0.95), Inches(0.92)])
    rect(slide, Inches(7.3), Inches(1.5), Inches(5.5), Inches(4.55), WHITE, LINE)
    textbox(slide, Inches(7.5), Inches(1.64), Inches(5.1), Inches(0.34), "Stocks organized per dataset", size=13, bold=True, color=INK)
    bar_rows = [(BAR_LABELS[i], int(DATA["targeted"][name]["total"][1])) for i, name in enumerate(TARGETED_ORDER)]
    hbars(slide, Inches(7.55), Inches(2.2), Inches(5.05), bar_rows, TEAL, rh=Inches(0.6), label_w=Inches(1.55))
    textbox(slide, Inches(0.55), Inches(6.2), Inches(12.2), Inches(0.34), "Counts are de-duplicated within each run (a stock is counted once across the six buckets).", size=9.5, color=MUT)
    footer(slide, 13)


def slide_summary_tx(prs):
    slide = blank_slide(prs)
    fill_bg(slide)
    header(slide, "Summary counts", "Transcriptomics run — totals")
    ribbon = [("16", "gene-set tiers", AMBER), ("5", "dataset categories", TEAL), (str(DATA["tx_total_stocks"]), "stocks organized", TEAL)]
    x = Inches(0.55)
    for val, lab, col in ribbon:
        stat_card(slide, x, Inches(1.5), val, lab, col)
        x += Inches(2.75)
    rect(slide, Inches(0.55), Inches(2.75), Inches(6.05), Inches(3.7), WHITE, LINE)
    textbox(slide, Inches(0.78), Inches(2.89), Inches(5.6), Inches(0.32), "By dataset category", size=12.5, bold=True, color=INK)
    tx_rows = [[TX_NM.get(r[0], r[0]), str(r[1]), str(r[2])] for r in DATA["tx_per_dataset"]]
    add_table(slide, Inches(0.78), Inches(3.26), Inches(5.6), ["Dataset category", "Tiers", "Stocks"], tx_rows, [Inches(3.5), Inches(0.95), Inches(1.15)])
    rect(slide, Inches(6.8), Inches(2.75), Inches(6.0), Inches(3.7), WHITE, LINE)
    textbox(slide, Inches(7.0), Inches(2.89), Inches(5.6), Inches(0.32), "By priority bucket (all 16 tiers)", size=12.5, bold=True, color=INK)
    b_rows = [(BUCKET_SHORT[i], int(b[1])) for i, b in enumerate(DATA["tx_by_bucket"])]
    hbars(slide, Inches(7.0), Inches(3.35), Inches(5.55), b_rows, AMBER, rh=Inches(0.48), label_w=Inches(1.75))
    textbox(slide, Inches(0.55), Inches(6.55), Inches(12.2), Inches(0.34), "Transcriptomics totals are summed across tiers; genes/alleles recur across tiers, so stocks are the headline metric.", size=9.5, color=MUT)
    footer(slide, 14)


def slide_examples(prs, n=9):
    slide = blank_slide(prs)
    fill_bg(slide)
    header(slide, "Worked examples", "One example reagent from each category")
    textbox(slide, Inches(0.55), Inches(1.34), Inches(12.2), Inches(0.32), "Representative rows from the GABA gene-list workbook (n=140) — every dataset follows the same six-table structure.", size=10.5, color=MUT)
    cw, ch, gx, gy = Inches(6.05), Inches(1.7), Inches(0.18), Inches(0.1)
    start_y = Inches(1.74)
    for i, ex in enumerate(DATA["examples"]):
        col, row = i % 2, i // 2
        X = Inches(0.55) + col * (cw + gx)
        Y = start_y + row * (ch + gy)
        rect(slide, X, Y, Inches(0.14), ch, BC[i])
        rect(slide, X + Inches(0.14), Y, cw - Inches(0.14), ch, WHITE, LINE)
        textbox(slide, X + Inches(0.32), Y + Inches(0.1), cw - Inches(0.5), Inches(0.3), ex["label"], size=11, bold=True, color=BC[i])
        textbox(slide, X + Inches(0.32), Y + Inches(0.4), cw - Inches(0.5), Inches(0.24), f"{ex['rows']} stocks in this bucket", size=8.5, color=MUT)
        for j, (lab, key) in enumerate([("ID:", "id"), ("Collection:", "collection"), ("Gene:", "gene"), ("Genotype:", "genotype")]):
            textbox(slide, X + Inches(0.32), Y + Inches(0.64) + Inches(0.24) * j, cw - Inches(0.55), Inches(0.24), f"{lab}  {ex[key]}", size=9, color=TX)
    footer(slide, n)


def slide_targeted_drilldown(prs, n=11):
    slide = blank_slide(prs)
    fill_bg(slide)
    header(slide, "Drill-down", "How targeted inputs differ before stock expansion")
    rows = [
        ["GABA/vGAT", "Fly gene list", "Gene, primary_symbol, FBgn, source dataset", "Direct stock expansion"],
        ["Bipolar GWAS", "Human genes -> DIOPT", "Human symbol, fly ortholog, FBgn, DIOPT score/rank, support tools", "Preserve orthology confidence"],
        ["AD GWAS", "Human genes -> DIOPT", "Human symbol, fly ortholog, FBgn, DIOPT rank, BDSC RNAi hints", "Use disease orthologs as candidate fly genes"],
        ["Jordan/Slumber", "Curated fly lists", "Submitted symbol, current symbol, FBgn, corrections", "Focused candidate follow-up"],
    ]
    add_table(
        slide,
        Inches(0.55),
        Inches(1.55),
        Inches(12.25),
        ["Input family", "Biological unit", "Important fields retained", "Why it matters"],
        rows,
        [Inches(1.75), Inches(2.05), Inches(5.35), Inches(3.1)],
        row_size=9,
        header_size=10,
    )
    callout(slide, Inches(0.55), Inches(4.05), Inches(5.95), Inches(1.45), "Why provenance comes first", "The first question is source and confidence. Column-level schema is useful only after the biological meaning and provenance are clear.", TEAL, body_size=10)
    callout(slide, Inches(6.85), Inches(4.05), Inches(5.95), Inches(1.45), "What remains auditable", "The generated workbooks still retain the detailed source columns and stock-level FlyBase fields needed to trace a reagent back to the input evidence.", AMBER, body_size=10)
    footer(slide, n)


def slide_tx_csw_drilldown(prs, n=12):
    slide = blank_slide(prs)
    fill_bg(slide)
    header(slide, "Drill-down", "Transcriptomics: consistent sleep/wake response tiers")
    callout(slide, Inches(0.55), Inches(1.45), Inches(5.95), Inches(1.28), "CSW FC1", "Higher-confidence, stricter log2 fold-change threshold (>1). Captures genes repeatedly changing with sleep/wake state across datasets.", TEAL, body_size=10)
    callout(slide, Inches(6.85), Inches(1.45), Inches(5.95), Inches(1.28), "CSW FC0.5", "Broader recurrence tier using a lower log2 fold-change threshold (>0.5) while preserving statistical rigor; useful for distributed homeostatic pathways.", AMBER, body_size=10)
    tiers = tier_rows(["CSW FC1", "CSW FC0.5"])
    add_table(slide, Inches(0.55), Inches(3.15), Inches(5.95), ["CSW tier", "Genes"], tiers[:8], [Inches(4.85), Inches(1.1)], row_size=8.5, header_size=9.5)
    add_table(slide, Inches(6.85), Inches(3.15), Inches(5.95), ["CSW tier", "Genes"], tiers[8:], [Inches(4.85), Inches(1.1)], row_size=8.5, header_size=9.5)
    textbox(slide, Inches(0.55), Inches(6.55), Inches(12.2), Inches(0.32), "Frequency = number of datasets showing consistent sleep/wake direction; higher frequency means more reproducible transcriptomic support.", size=9.5, color=MUT)
    footer(slide, n)


def slide_tx_correlation_drilldown(prs, n=13):
    slide = blank_slide(prs)
    fill_bg(slide)
    header(slide, "Drill-down", "Transcriptomics: history, rebound, and literature support")
    rows = [
        ["SleepHistory correlates", "215 genes", "Genes whose expression correlates with prior sleep/wake history across time bins."],
        ["History x Rebound overlaps", "20 genes", "Genes linking prior history direction with predicted rebound direction."],
        ["Homeostatic genes", "6 genes", "Small focused set called out as sleep-homeostasis candidates."],
        ["fl.ai annotations", "per gene", "AI-assisted literature category, confidence, rationale, supporting references, and source papers."],
    ]
    add_table(slide, Inches(0.55), Inches(1.55), Inches(12.25), ["Evidence layer", "Size", "Interpretation"], rows, [Inches(3.0), Inches(1.4), Inches(7.85)], row_size=9.5, header_size=10.5)
    callout(slide, Inches(0.55), Inches(4.0), Inches(5.95), Inches(1.55), "Correlation logic", "The paper asks whether expression tracks sleep-wake history before sampling and/or predicts rebound after sampling, using Pearson correlations against log2-normalized gene expression.", TEAL, body_size=9.5)
    callout(slide, Inches(6.85), Inches(4.0), Inches(5.95), Inches(1.55), "How this affects stock choice", "A gene that is both transcriptomically recurrent and literature-supported is easier to justify for reagent ordering than a single-method hit with no sleep/circadian evidence.", AMBER, body_size=9.5)
    footer(slide, n)


def slide_decision_use(prs, n=14):
    slide = blank_slide(prs)
    fill_bg(slide)
    header(slide, "How to use the output", "A fast decision path")
    rows = [
        ("Start with source confidence", "Was this gene from direct fly evidence, a human GWAS ortholog, or a transcriptomics recurrence/correlation tier?"),
        ("Check orderability", "Prefer Bloomington/Vienna stock-center reagents before non-stock-center or custom reagents."),
        ("Check reagent type", "Separate UAS/RNAi perturbation from allele/insertion reagents depending on the experiment."),
        ("Check biological support", "Use reference counts, sleep/circadian flags, phenotype relevance, and fl.ai rationale to prioritize follow-up."),
        ("Shortlist manually", "Open the relevant workbook sheet and review genotype/construct details before ordering."),
    ]
    bullet_rows(slide, Inches(0.8), Inches(1.55), Inches(11.7), rows, size=12, gap=0.78)
    footer(slide, n)


def slide_final_summary_counts(prs, n=15):
    slide = blank_slide(prs)
    fill_bg(slide)
    header(slide, "Final summary", "Where the strongest reagent pools are")
    bar_rows = [(BAR_LABELS[i], int(DATA["targeted"][name]["total"][1])) for i, name in enumerate(TARGETED_ORDER)]
    rect(slide, Inches(0.55), Inches(1.5), Inches(5.9), Inches(4.85), WHITE, LINE)
    textbox(slide, Inches(0.78), Inches(1.68), Inches(5.4), Inches(0.32), "Targeted runs: stocks per dataset", size=13, bold=True, color=INK)
    hbars(slide, Inches(0.75), Inches(2.25), Inches(5.35), bar_rows, TEAL, rh=Inches(0.55), label_w=Inches(1.55))
    rect(slide, Inches(6.85), Inches(1.5), Inches(5.95), Inches(4.85), WHITE, LINE)
    textbox(slide, Inches(7.08), Inches(1.68), Inches(5.4), Inches(0.32), "Transcriptomics: stocks per bucket", size=13, bold=True, color=INK)
    b_rows = [(BUCKET_SHORT[i], int(b[1])) for i, b in enumerate(DATA["tx_by_bucket"])]
    hbars(slide, Inches(7.05), Inches(2.25), Inches(5.35), b_rows, AMBER, rh=Inches(0.48), label_w=Inches(1.75))
    textbox(slide, Inches(0.55), Inches(6.55), Inches(12.2), Inches(0.34), "This is the slide to use for a quick sense of scale before opening individual workbooks.", size=9.5, color=MUT)
    footer(slide, n)


def slide_closing(prs):
    slide = blank_slide(prs)
    fill_bg(slide, INK)
    rect(slide, Inches(0), Inches(0), Inches(0.28), H, TEAL)
    rect(slide, Inches(0.28), Inches(0), Inches(0.08), H, AMBER)
    textbox(slide, Inches(1), Inches(1), Inches(11), Inches(0.4), "TAKEAWAYS", size=13, bold=True, color=TEALL)
    textbox(slide, Inches(0.95), Inches(1.4), Inches(11.6), Inches(0.8), "What the tables give you", size=34, bold=True, color=WHITE, font="Georgia")
    pts = [
        ("22 gene sets, one consistent format", "Every input list becomes the same six prioritized, phenotype-filtered stock tables."),
        ("Ordering by orderability + reagent type", "Bloomington / Vienna / other centers first; UAS vs classical alleles separated; custom reagents last."),
        ("Built-in literature relevance", "Sleep / circadian / rhythm / locomotor references and phenotype-similarity scores travel with each stock."),
        ("8,300+ phenotype-backed stocks organized", "Ready to filter down to an orderable shortlist per gene."),
    ]
    y = Inches(2.7)
    for i, (t, d) in enumerate(pts):
        oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1), y, Inches(0.42), Inches(0.42))
        oval.fill.solid()
        oval.fill.fore_color.rgb = TEAL
        oval.line.fill.background()
        textbox(slide, Inches(1), y, Inches(0.42), Inches(0.42), str(i + 1), size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        textbox(slide, Inches(1.6), y - Inches(0.04), Inches(11), Inches(0.36), t, size=15, bold=True, color=WHITE)
        textbox(slide, Inches(1.6), y + Inches(0.32), Inches(11), Inches(0.4), d, size=11.5, color=LIGHT)
        y += Inches(0.95)
    footer(slide, 16)


def main():
    prs = new_prs()
    slide_title(prs)
    slide_executive_summary(prs)
    slide_pipeline(prs)
    slide_source_map(prs)
    slide_targeted_provenance(prs)
    slide_transcriptomics_provenance(prs)
    slide_counts_overview(prs)
    slide_med24_example(prs)
    slide_examples(prs, 9)
    slide_anatomy(prs, 10)
    slide_targeted_drilldown(prs, 11)
    slide_tx_csw_drilldown(prs, 12)
    slide_tx_correlation_drilldown(prs, 13)
    slide_decision_use(prs, 14)
    slide_final_summary_counts(prs, 15)
    slide_closing(prs)

    prs.core_properties.title = "FlyBase Reagent Stocker — Pipeline Overview"
    prs.core_properties.author = "Allada Lab"
    prs.save(str(OUT))

    from sanitize_pptx import sanitize

    sanitize(OUT)
    print(f"WROTE {OUT}")


if __name__ == "__main__":
    main()
