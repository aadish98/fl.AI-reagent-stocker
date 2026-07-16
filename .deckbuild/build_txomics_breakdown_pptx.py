#!/usr/bin/env python3
"""Build Tx-Omics gene-set overview deck with python-pptx (PowerPoint-native OOXML)."""
from __future__ import annotations

import csv
import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).parent
BREAKDOWN = ROOT.parent / "data/gene_sets/Tx-Omics-FollowUp_v3/Breakdown"
TEMPLATE = ROOT.parent / "Allada_Stocker_Tables_Overview.pptx"
OUT = BREAKDOWN / "Tx-Omics_GeneSets_Overview.pptx"
TOTAL = 4

# Mac PowerPoint repair is triggered by some Unicode punctuation in OOXML text runs.
_ASCII_REPLACEMENTS = str.maketrans({
    "\u201c": '"', "\u201d": '"', "\u2018": "'", "\u2019": "'",
    "\u2014": "-", "\u2212": "-", "\u2192": "->", "\u2082": "2",
    "\u00b7": " / ", "\u2229": " & ",
})


def safe(text: str) -> str:
    return text.translate(_ASCII_REPLACEMENTS)

W = Inches(13.333)
H = Inches(7.5)

INK = RGBColor(0x15, 0x23, 0x3B)
TEAL = RGBColor(0x0E, 0x7C, 0x7B)
TEALL = RGBColor(0x8F, 0xCF, 0xCD)
AMBER = RGBColor(0xE0, 0xA1, 0x3A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG = RGBColor(0xF4, 0xF6, 0xF8)
TX = RGBColor(0x1F, 0x2A, 0x37)
MUT = RGBColor(0x66, 0x73, 0x85)
LINE = RGBColor(0xE2, 0xE8, 0xF0)
PALE = RGBColor(0xEE, 0xF2, 0xF5)
LIGHT = RGBColor(0xC7, 0xD2, 0xDE)
BLUE = RGBColor(0x1C, 0x6F, 0xB0)

BC = [TEAL, BLUE, RGBColor(0x6B, 0x4F, 0xA0), AMBER, RGBColor(0xB2, 0x41, 0x5A), RGBColor(0x4F, 0x7A, 0x2E)]


def new_prs() -> Presentation:
    if TEMPLATE.exists():
        prs = Presentation(str(TEMPLATE))
        sld_id_lst = prs.slides._sldIdLst
        for sld_id in list(sld_id_lst):
            r_id = sld_id.rId
            prs.part.drop_rel(r_id)
            sld_id_lst.remove(sld_id)
        return prs
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    return prs


def blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def fill_bg(slide, color=BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def rect(slide, left, top, width, height, fill, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line:
        sh.line.color.rgb = line
        sh.line.width = Pt(0.75)
    else:
        sh.line.fill.background()
    return sh


def tb(slide, left, top, width, height, text, *, size=11, bold=False, color=TX, font="Calibri", align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = safe(text)
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font
    p.font.color.rgb = color
    p.alignment = align
    return box


def header(slide, kicker, title):
    rect(slide, Inches(0), Inches(0), W, Inches(1.18), INK)
    rect(slide, Inches(0), Inches(1.18), W, Inches(0.06), AMBER)
    if kicker:
        tb(slide, Inches(0.55), Inches(0.18), Inches(12), Inches(0.3), kicker.upper(), size=11, bold=True, color=TEALL)
    tb(slide, Inches(0.55), Inches(0.44), Inches(12.2), Inches(0.62), title, size=25, bold=True, color=WHITE, font="Georgia")


def footer(slide, n):
    rect(slide, Inches(0), Inches(7.18), W, Inches(0.32), INK)
    tb(slide, Inches(0.5), Inches(7.16), Inches(8), Inches(0.32), "Allada Lab  ·  Tx-Omics Follow-Up v3", size=9, color=TEALL)
    tb(slide, Inches(11.8), Inches(7.16), Inches(1.2), Inches(0.32), f"{n} / {TOTAL}", size=9, color=TEALL, align=PP_ALIGN.RIGHT)


def card(slide, left, top, width, height):
    rect(slide, left, top, width, height, WHITE, LINE)


def style_cell(cell, text, *, bold=False, size=10, color=TX, fill=None, align=PP_ALIGN.LEFT):
    cell.text = safe(str(text))
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if fill:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.name = "Calibri"
        p.font.color.rgb = color
        p.alignment = align


def add_table(slide, left, top, width, headers, rows, col_widths, *, row_h=0.32, body_size=10, header_size=10.5):
    tbl = slide.shapes.add_table(len(rows) + 1, len(headers), left, top, width, Inches(row_h * (len(rows) + 1))).table
    for i, cw in enumerate(col_widths):
        tbl.columns[i].width = cw
    for j, h in enumerate(headers):
        style_cell(tbl.cell(0, j), h, bold=True, size=header_size, color=WHITE, fill=INK)
    for i, row in enumerate(rows):
        bg = PALE if i % 2 else WHITE
        for j, val in enumerate(row):
            if isinstance(val, tuple):
                txt, bold, col = val
                style_cell(tbl.cell(i + 1, j), txt, bold=bold, size=body_size, color=col, fill=bg)
            else:
                style_cell(tbl.cell(i + 1, j), val, size=body_size, fill=bg)
    return tbl


def stat_big(slide, left, top, width, value, label, color):
    card(slide, left, top, width, Inches(2.35))
    tb(slide, left, top + Inches(0.2), width, Inches(0.65), str(value), size=30, bold=True, color=color, font="Georgia", align=PP_ALIGN.CENTER)
    tb(slide, left, top + Inches(0.85), width, Inches(0.35), label, size=10.5, color=MUT, align=PP_ALIGN.CENTER)


def hbars(slide, left, top, width, rows, color, *, rh=0.42, label_w=1.35, val_w=0.55):
    max_v = max(v for _, v in rows) or 1
    bar_max = width - Inches(label_w) - Inches(val_w) - Inches(0.12)
    for i, (lab, val) in enumerate(rows):
        y = top + Inches(rh * i)
        tb(slide, left, y, Inches(label_w - 0.1), Inches(rh), lab, size=9, color=MUT, align=PP_ALIGN.RIGHT)
        bw = max(Inches(0.05), int(bar_max * val / max_v))
        rect(slide, left + Inches(label_w), y + Inches(rh * 0.2), bw, Inches(rh * 0.6), color)
        tb(slide, left + Inches(label_w) + bw + Inches(0.06), y, Inches(val_w), Inches(rh), str(val), size=9, bold=True, color=TX)


def bullets(slide, left, top, width, items, *, size=11, color=LIGHT):
    box = slide.shapes.add_textbox(left, top, width, Inches(3.5))
    tf = box.text_frame
    tf.word_wrap = True
    for i, text in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = safe(f"• {text}")
        p.level = 0
        p.font.size = Pt(size)
        p.font.name = "Calibri"
        p.font.color.rgb = color


def parse_gene_sets():
    combo: dict[str, dict] = {}
    with (BREAKDOWN / "combination_counts_summary.csv").open(newline="", encoding="utf-8") as fh:
        for row in csv.DictReader(fh):
            gs = row["gene_set"]
            combo.setdefault(gs, {"genes": 0, "stocks": 0, "alleles": 0})
            combo[gs]["genes"] += int(row["num_genes"])
            combo[gs]["stocks"] += int(row["num_stocks"])
            combo[gs]["alleles"] += int(row["num_alleles"])

    sets = []
    for path in sorted(BREAKDOWN.glob("*.csv")):
        if path.name == "combination_counts_summary.csv":
            continue
        name = path.stem
        m = re.search(r"n=(\d+)genes", name)
        genes = int(m.group(1)) if m else sum(1 for _ in path.open()) - 1
        if name.startswith("FC0.5_Sleep"):
            cat = "FC0.5 · Sleep rebound"
        elif name.startswith("FC0.5_Wake"):
            cat = "FC0.5 · Wake rebound"
        elif name.startswith("FC1_Sleep"):
            cat = "FC1 · Sleep rebound"
        elif name.startswith("FC1_Wake"):
            cat = "FC1 · Wake rebound"
        elif "History" in name and "overlap" in name:
            cat = "History × Rebound overlap"
        elif "History" in name:
            cat = "Sleep history correlates"
        elif "homeostatic" in name:
            cat = "Curated · Homeostatic core"
        elif "Mechanistic" in name:
            cat = "Curated · Mechanistic screen"
        else:
            cat = "Other"
        c = combo.get(name, {})
        sets.append({"name": name, "genes": genes, "cat": cat, "stocks": c.get("stocks", 0), "alleles": c.get("alleles", 0)})
    sets.sort(key=lambda x: x["genes"], reverse=True)

    mech = []
    with (BREAKDOWN / "Tx-Omics_Mechanistic_Screen_Candidates_nsyb_elav.csv").open(newline="", encoding="utf-8-sig") as fh:
        reader = csv.DictReader(fh)
        for row in reader:
            mech.append({
                "rank": row["Priority Rank"],
                "gene": row["Gene"],
                "category": row["Mechanistic Category"].split(" / ")[0],
                "priority": row["Screen Priority"],
            })
    return {"sets": sets, "mech": mech}


def short_name(name: str) -> str:
    return (
        name.replace(re.search(r"_n=\d+genes", name).group(0) if re.search(r"_n=\d+genes", name) else "", "")
        .replace(" Table", "")
        .replace("frequency_", "freq ")
        .replace("_", " ")
        .replace("Tx-Omics ", "")
    )


def slide_title(prs, data):
    s = blank(prs)
    fill_bg(s, INK)
    rect(s, Inches(0), Inches(0), Inches(0.28), H, TEAL)
    rect(s, Inches(0.28), Inches(0), Inches(0.08), H, AMBER)
    tb(s, Inches(1.0), Inches(1.5), Inches(11), Inches(0.38), "TX-OMICS FOLLOW-UP v3", size=13, bold=True, color=TEALL)
    tb(s, Inches(0.95), Inches(2.0), Inches(11.6), Inches(1.0), "Gene Set Checkpoint", size=44, bold=True, color=WHITE, font="Georgia")
    tb(s, Inches(1.0), Inches(3.1), Inches(11.2), Inches(0.45), "Quick reference for follow-up gene lists and screening priorities", size=17, color=LIGHT)
    tb(s, Inches(1.0), Inches(3.6), Inches(11), Inches(0.35), "Rosensweig, Shah et al. 2026", size=12, italic=True, color=TEALL)
    stats = [
        ("17", "gene lists"),
        ("1,162", "genes total"),
        (f'{sum(x["stocks"] for x in data["sets"]):,}', "stocks ready"),
        ("25", "priority candidates"),
    ]
    x = Inches(1.0)
    for val, lab in stats:
        tb(s, x, Inches(4.6), Inches(2.65), Inches(0.58), val, size=28, bold=True, color=AMBER, font="Georgia")
        tb(s, x, Inches(5.18), Inches(2.65), Inches(0.45), lab, size=11, color=LIGHT)
        x += Inches(2.85)
    tb(s, Inches(1.0), Inches(6.55), Inches(8), Inches(0.35), "Allada Lab  ·  July 2026", size=12, italic=True, color=TEALL)


def slide_map(prs, data):
    """Four buckets - what exists and how big each is."""
    s = blank(prs)
    fill_bg(s)
    header(s, "At a glance", "Four types of gene lists")

    def count_cat(prefix):
        items = [x for x in data["sets"] if x["cat"].startswith(prefix) or prefix in x["cat"]]
        return len(items), sum(x["genes"] for x in items)

    rebound_n, rebound_g = count_cat("FC")
    hist = [x for x in data["sets"] if x["cat"] == "Sleep history correlates"]
    ov = [x for x in data["sets"] if "overlap" in x["cat"]]
    curated = [x for x in data["sets"] if x["cat"].startswith("Curated")]

    rows = [
        ["Rebound responders", f"{rebound_n} lists", f"{rebound_g} genes", "Genes that change consistently after sleep deprivation across manipulations"],
        ["Sleep history correlates", f"{len(hist)} lists", f"{sum(x['genes'] for x in hist)} genes", "Expression tracks prior sleep or wake history"],
        ["History / rebound overlaps", f"{len(ov)} lists", f"{sum(x['genes'] for x in ov)} genes", "Dual signature - strongest homeostatic candidates"],
        ["Curated screen list", f"{len(curated)} lists", f"{sum(x['genes'] for x in curated)} genes", "6 core homeostatic genes + 19 ranked for RNAi screen"],
    ]
    card(s, Inches(0.55), Inches(1.5), Inches(12.2), Inches(4.8))
    add_table(s, Inches(0.78), Inches(1.65), Inches(11.7), ["Category", "Lists", "Genes", "What it means"], rows,
              [Inches(2.4), Inches(0.9), Inches(1.0), Inches(7.4)], row_h=0.55, body_size=11, header_size=11)

    tb(s, Inches(0.55), Inches(6.45), Inches(12.2), Inches(0.45),
       "Wake-enriched rebound lists are the largest. Overlap and curated lists are smallest but highest priority for mechanistic follow-up.",
       size=11, italic=True, color=MUT)
    footer(s, 2)


def slide_start_here(prs, data):
    """PI action slide - where to focus first."""
    s = blank(prs)
    fill_bg(s)
    header(s, "Start here", "Priority genes for first-pass screening")

    card(s, Inches(0.55), Inches(1.5), Inches(5.9), Inches(5.0))
    tb(s, Inches(0.78), Inches(1.65), Inches(5.4), Inches(0.35), "Homeostatic core (6 genes)", size=14, bold=True, color=INK)
    tb(s, Inches(0.78), Inches(2.05), Inches(5.4), Inches(0.5),
       "unc79  /  SIFa  /  rumpel  /  AstA-R2  /  Trhn  /  RpL23", size=12, bold=True, color=TEAL)
    tb(s, Inches(0.78), Inches(2.65), Inches(5.4), Inches(1.0),
       "NALCN channel, neuropeptide signaling, glial transport, serotonin biosynthesis. Literature-supported sleep homeostasis.",
       size=11, color=MUT)

    card(s, Inches(0.78), Inches(3.85), Inches(5.4), Inches(2.4))
    tb(s, Inches(0.98), Inches(4.0), Inches(5.0), Inches(0.3), "Top screen candidates", size=12, bold=True, color=INK)
    top_genes = ", ".join(m["gene"] for m in data["mech"][:8])
    tb(s, Inches(0.98), Inches(4.35), Inches(5.0), Inches(0.55), top_genes, size=10.5, color=TEAL)
    tb(s, Inches(0.98), Inches(4.95), Inches(5.0), Inches(1.0),
       "Ranked list of 19 genes for a first-pass neuronal knockdown screen. Includes na, unc80, Mip, BomBc2.",
       size=10, color=MUT)

    card(s, Inches(6.65), Inches(1.5), Inches(6.1), Inches(5.0))
    tb(s, Inches(6.88), Inches(1.65), Inches(5.6), Inches(0.35), "Why these first?", size=14, bold=True, color=INK)
    bullets(s, Inches(6.88), Inches(2.1), Inches(5.6), [
        "Overlap genes (20 total) show opposing history vs. rebound dynamics - prime regulators.",
        "Homeostatic core converges on NALCN, neuropeptide, and glial pathways from the paper.",
        "Reagents are already mapped for every list - ready to order.",
    ], size=12, color=TX)

    tb(s, Inches(6.88), Inches(4.5), Inches(5.6), Inches(0.35), "Overlap highlights", size=12, bold=True, color=INK)
    bullets(s, Inches(6.88), Inches(4.85), Inches(5.6), [
        "History(-) x Rebound(+) : 5 genes (e.g. CG9377)",
        "History(+) x Rebound(-) : 15 genes (e.g. AstA-R2, Mip)",
    ], size=11, color=MUT)

    footer(s, 3)


def slide_where(prs):
    """Simple navigation - no file paths."""
    s = blank(prs)
    fill_bg(s, INK)
    rect(s, Inches(0), Inches(7.44), W, Inches(0.06), AMBER)
    tb(s, Inches(0.75), Inches(1.3), Inches(11), Inches(0.55), "Where to find everything", size=32, bold=True, color=WHITE, font="Georgia")

    items = [
        ("Gene lists", "Tx-Omics Follow-Up v3 folder - one list per file, genes annotated with literature and cycling status."),
        ("Stock tables", "Stock tables folder - Bloomington and VDRC lines, alleles, prioritized by availability."),
        ("Summary counts", "Summary spreadsheet - stocks and alleles per list at a glance."),
    ]
    y = Inches(2.2)
    for i, (title, body) in enumerate(items):
        rect(s, Inches(0.75), y, Inches(0.1), Inches(1.15), BC[i])
        card(s, Inches(0.85), y, Inches(11.7), Inches(1.15))
        tb(s, Inches(1.05), y + Inches(0.12), Inches(3.0), Inches(0.35), title, size=14, bold=True, color=AMBER)
        tb(s, Inches(1.05), y + Inches(0.48), Inches(11.2), Inches(0.55), body, size=12, color=LIGHT)
        y += Inches(1.35)

    tb(s, Inches(0.75), Inches(6.3), Inches(11.5), Inches(0.5),
       "Suggested next step: order reagents for overlap + homeostatic core, then run neuronal knockdown sleep assays.",
       size=13, color=TEALL)
    footer(s, 4)


def powerpoint_native_save(path: Path) -> bool:
    """Round-trip through Microsoft PowerPoint so Mac PowerPoint opens without repair."""
    import subprocess
    import tempfile

    tmp = Path(tempfile.gettempdir()) / "txomics_pptx_native.pptx"
    if tmp.exists():
        tmp.unlink()
    script = f'''
set inPath to POSIX file "{path}"
set outPath to POSIX file "{tmp}"

tell application "Microsoft PowerPoint"
    activate
    open inPath
    delay 6
    save active presentation in outPath
    close active presentation saving no
end tell
'''
    try:
        subprocess.run(["osascript", "-e", 'tell application "Microsoft PowerPoint" to quit'], check=False, capture_output=True)
        subprocess.run(["sleep", "1"], check=False)
        result = subprocess.run(["osascript", "-e", script], check=False, capture_output=True, text=True)
        if result.returncode != 0 or not tmp.exists():
            print("WARN: PowerPoint native save skipped:", (result.stderr or result.stdout or "").strip())
            return False
        tmp.replace(path)
        subprocess.run(["xattr", "-cr", str(path)], check=False)
        return True
    except Exception as exc:
        print(f"WARN: PowerPoint native save failed: {exc}")
        return False


def main():
    import subprocess

    data = parse_gene_sets()
    prs = new_prs()
    slide_title(prs, data)
    slide_map(prs, data)
    slide_start_here(prs, data)
    slide_where(prs)
    prs.core_properties.title = "Tx-Omics Follow-Up v3 — Gene Set Checkpoint"
    prs.core_properties.author = "Allada Lab"
    prs.save(str(OUT))

    from sanitize_pptx import sanitize

    sanitize(OUT)
    if powerpoint_native_save(OUT):
        print(f"WROTE {OUT} (PowerPoint-native)")
    else:
        subprocess.run(["xattr", "-cr", str(OUT)], check=False)
        print(f"WROTE {OUT}")


if __name__ == "__main__":
    main()
